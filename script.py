import os
import google.generativeai as genai
import pdfplumber
import streamlit as st
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation
from PIL import Image
from io import BytesIO
from dotenv import load_dotenv
import pygame
import random

# Load environment variables from .env file
load_dotenv()

# Configure the Gemini API
genai.configure(api_key=os.environ["GEMINI_API_KEY"])

# Set up generation configuration for the model
generation_config = {
    "temperature": 1,
    "top_p": 0.95,
    "top_k": 40,
    "max_output_tokens": 8192,
    "response_mime_type": "text/plain",
}

# Initialize the model with the configuration
model = genai.GenerativeModel(
    model_name="gemini-1.5-pro-002",
    generation_config=generation_config,
)

# Helper functions to extract text from different file types
def extract_text_from_pdf(file):
    return extract_text(file)

def extract_text_from_docx(file):
    doc = Document(file)
    return '\n'.join([paragraph.text for paragraph in doc.paragraphs])

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

def extract_text_from_image(file):
    image = Image.open(BytesIO(file.read()))
    ocr_text = genai.ocr(image)  # Placeholder for Gemini OCR or external OCR function
    return ocr_text

# Function to perform a task on the extracted text using Gemini
def perform_task(task_prompt, content):
    chat_session = model.start_chat(history=[])
    response = chat_session.send_message(f"{task_prompt}\n{content}")
    return response.text

# Game functions
def start_pygame():
    pygame.init()
    WIDTH, HEIGHT = 600, 600
    win = pygame.display.set_mode((WIDTH, HEIGHT))
    pygame.display.set_caption("Game Selector")

    WHITE = (255, 255, 255)
    BLACK = (0, 0, 0)
    BLUE = (0, 0, 255)
    GREEN = (0, 255, 0)
    RED = (255, 0, 0)

    FONT = pygame.font.SysFont('Arial', 40)

    def draw_text(text, font, color, surface, x, y):
        textobj = font.render(text, True, color)
        textrect = textobj.get_rect()
        textrect.topleft = (x, y)
        surface.blit(textobj, textrect)

    def main_menu():
        run = True
        while run:
            win.fill(WHITE)
            draw_text('Press SPACE to Play Tic Tac Toe', FONT, BLACK, win, 20, 20)
            draw_text('Press ENTER to Play Rock Paper Scissors', FONT, BLACK, win, 20, 100)
            pygame.display.update()

            for event in pygame.event.get():
                if event.type == pygame.QUIT:
                    run = False
                if event.type == pygame.KEYDOWN:
                    if event.key == pygame.K_SPACE:
                        tic_tac_toe()
                    if event.key == pygame.K_RETURN:
                        rock_paper_scissors()

    def tic_tac_toe():
        board = [' ' for _ in range(9)]
        player = 'X'
        ai = 'O'
        game_over = False

        def draw_board():
            win.fill(WHITE)
            for i in range(3):
                pygame.draw.line(win, BLACK, (0, (i + 1) * HEIGHT // 3), (WIDTH, (i + 1) * HEIGHT // 3), 5)
                pygame.draw.line(win, BLACK, ((i + 1) * WIDTH // 3, 0), ((i + 1) * WIDTH // 3, HEIGHT), 5)

            for i in range(9):
                if board[i] == 'X':
                    draw_text('X', FONT, RED, win, (i % 3) * (WIDTH // 3) + WIDTH // 6 - 20, (i // 3) * (HEIGHT // 3) + HEIGHT // 6 - 20)
                elif board[i] == 'O':
                    draw_text('O', FONT, GREEN, win, (i % 3) * (WIDTH // 3) + WIDTH // 6 - 20, (i // 3) * (HEIGHT // 3) + HEIGHT // 6 - 20)

            pygame.display.update()

        def check_winner():
            win_conditions = [(0, 1, 2), (3, 4, 5), (6, 7, 8),
                              (0, 3, 6), (1, 4, 7), (2, 5, 8),
                              (0, 4, 8), (2, 4, 6)]
            for condition in win_conditions:
                if board[condition[0]] == board[condition[1]] == board[condition[2]] != ' ':
                    return board[condition[0]]
            if ' ' not in board:
                return 'Tie'
            return None

        def ai_move():
            for i in range(9):
                if board[i] == ' ':
                    board[i] = ai
                    if check_winner() == ai:
                        return
                    board[i] = ' '

            for i in range(9):
                if board[i] == ' ':
                    board[i] = player
                    if check_winner() == player:
                        board[i] = ai
                        return
                    board[i] = ' '

            available_moves = [i for i in range(9) if board[i] == ' ']
            if available_moves:
                board[random.choice(available_moves)] = ai

        while not game_over:
            draw_board()

            for event in pygame.event.get():
                if event.type == pygame.QUIT:
                    game_over = True
                if event.type == pygame.MOUSEBUTTONDOWN and player == 'X':
                    pos = pygame.mouse.get_pos()
                    index = (pos[1] // (HEIGHT // 3)) * 3 + (pos[0] // (WIDTH // 3))
                    if board[index] == ' ':
                        board[index] = player
                        winner = check_winner()
                        if winner is not None:
                            draw_board()
                            if winner == 'Tie':
                                draw_text('Tie!', FONT, BLACK, win, 20, 20)
                            else:
                                draw_text(f'{winner} wins!', FONT, BLACK, win, 20, 20)
                            pygame.display.update()
                            pygame.time.delay(2000)
                            return

                        ai_move()
                        winner = check_winner()
                        if winner is not None:
                            draw_board()
                            if winner == 'Tie':
                                draw_text('Tie!', FONT, BLACK, win, 20, 20)
                            else:
                                draw_text(f'{winner} wins!', FONT, BLACK, win, 20, 20)
                            pygame.display.update()
                            pygame.time.delay(2000)
                            return

                        player = 'X'

    def rock_paper_scissors():
        choices = ['Rock', 'Paper', 'Scissors']
        run = True
        while run:
            win.fill(WHITE)
            draw_text('Rock Paper Scissors', FONT, BLACK, win, 20, 20)
            draw_text('Press 1: Rock', FONT, BLUE, win, 20, 100)
            draw_text('Press 2: Paper', FONT, BLUE, win, 20, 200)
            draw_text('Press 3: Scissors', FONT, BLUE, win, 20, 300)
            pygame.display.update()

            for event in pygame.event.get():
                if event.type == pygame.QUIT:
                    run = False
                if event.type == pygame.KEYDOWN:
                    if event.key == pygame.K_1:
                        player_choice = 'Rock'
                    elif event.key == pygame.K_2:
                        player_choice = 'Paper'
                    elif event.key == pygame.K_3:
                        player_choice = 'Scissors'
                    else:
                        continue

                    ai_choice = random.choice(choices)
                    draw_text(f'You chose: {player_choice}', FONT, BLACK, win, 20, 400)
                    draw_text(f'AI chose: {ai_choice}', FONT, BLACK, win, 20, 450)

                    if player_choice == ai_choice:
                        draw_text('It\'s a Tie!', FONT, BLACK, win, 20, 500)
                    elif (player_choice == 'Rock' and ai_choice == 'Scissors') or \
                         (player_choice == 'Paper' and ai_choice == 'Rock') or \
                         (player_choice == 'Scissors' and ai_choice == 'Paper'):
                        draw_text('You Win!', FONT, BLACK, win, 20, 500)
                    else:
                        draw_text('You Lose!', FONT, BLACK, win, 20, 500)

                    pygame.display.update()
                    pygame.time.delay(2000)
                    return

    main_menu()
    pygame.quit()

# Streamlit app layout
st.sidebar.title("Menu")
st.sidebar.write("Upload your file and specify a task for processing.")
st.sidebar.write("Supported file types: PDF, DOCX, PPTX, JPG, JPEG, PNG")

# Display logo at the top (replace 'logo.png' with the path to your logo)
st.image("logo.jpg", width=200)

# Main app content
st.write("## Upload a File and Enter a Prompt")
st.write("This app extracts content from PDF, Word, PPT, or image files and performs a task based on your prompt.")

# File upload and task prompt input
uploaded_file = st.file_uploader("Choose a file", type=["pdf", "docx", "pptx", "jpg", "jpeg", "png"])
task_prompt = st.text_input("Enter your task prompt (e.g., 'Summarize', 'Generate flashcards', 'Create quiz', etc.)")

# Game button
if st.button("Play Game"):
    st.write("Launching game...")
    start_pygame()

# Check if a file and task prompt are provided
if uploaded_file and task_prompt:
    file_type = uploaded_file.name.split('.')[-1]

    # Extract content based on file type
    if file_type == "pdf":
        content = extract_text_from_pdf(uploaded_file)
    elif file_type == "docx":
        content = extract_text_from_docx(uploaded_file)
    elif file_type == "pptx":
        content = extract_text_from_pptx(uploaded_file)
    elif file_type in ["jpg", "jpeg", "png"]:
        content = extract_text_from_image(uploaded_file)
    else:
        st.write("Unsupported file type.")
        content = None

    # Perform the task if content is successfully extracted
    if content:
        st.write("### Task Result")
        result = perform_task(task_prompt, content)
        st.write(result)
    else:
        st.write("Could not extract content from the uploaded file.")
